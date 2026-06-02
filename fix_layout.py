import sys
import io
import math
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Pt, Emu

IN = 914400  # EMU per inch

def em(inches): return int(inches * IN)
def inch(emu): return emu / IN

STANDARD_PT = 14  # Unified font size

src = r'C:\Users\JasonOng\OneDrive - Advancer Global Ltd\AST BD\2024 HyESys\Software\Manual logic - bef AI\HZ.SG v2.1 logic2 (EN).pptx'
dst = r'C:\Users\JasonOng\OneDrive - Advancer Global Ltd\AST BD\2024 HyESys\Software\Manual logic - bef AI\HZ.SG v2.1 logic2 (EN).pptx'

prs = Presentation(src)
SW = prs.slide_width   # 6858000 EMU = 7.5"
SH = prs.slide_height  # 12192000 EMU = 13.333"

MARGIN = em(0.05)

def set_font(tf, pt):
    """Force all runs in a text frame to the given font size."""
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(pt)

def line_count_estimate(tf, box_w_emu):
    """Estimate visual line count for height calculation."""
    box_w_in = inch(box_w_emu)
    chars_per_line = max(8, int(box_w_in * 10))  # ~10 chars/inch at 14pt
    total = 0
    for para in tf.paragraphs:
        text = para.text.strip()
        if not text:
            total += 0.4
        else:
            total += max(1, math.ceil(len(text) / chars_per_line))
    return max(1.0, total)

def calc_height(tf, box_w_emu, pt=STANDARD_PT):
    n_lines = line_count_estimate(tf, box_w_emu)
    line_h = pt * 1.2 / 72  # inches per line
    return em(n_lines * line_h + 0.12)  # 0.12" top+bottom padding

def clamp(shape):
    """Nudge shape back into slide bounds without changing content size (for rectangles)."""
    L = shape.left or 0
    T = shape.top or 0
    W = shape.width or 0
    H = shape.height or 0
    if L < MARGIN: L = MARGIN
    if T < MARGIN: T = MARGIN
    if L + W > SW - MARGIN: L = SW - W - MARGIN
    if T + H > SH - MARGIN: T = SH - H - MARGIN
    shape.left, shape.top = L, T

def resize_right(shape):
    """Trim width so right edge stays within slide."""
    L = shape.left or 0
    W = shape.width or 0
    if L + W > SW - MARGIN:
        shape.width = max(em(0.5), SW - L - MARGIN)

def move_and_size(shape, left=None, top=None, width=None, height=None):
    if left  is not None: shape.left   = em(left)
    if top   is not None: shape.top    = em(top)
    if width is not None: shape.width  = em(width)
    if height is not None: shape.height = em(height)


# ── Per-slide corrections ─────────────────────────────────────────────────────

# Slide 2 (slide_idx=1): circuit diagram
def fix_slide2(slide):
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        tf = shape.text_frame
        set_font(tf, STANDARD_PT)
        nm = shape.name

        if nm == 'TextBox 29':   # "Load Meter 607" — nudge in from left
            move_and_size(shape, left=0.06, top=2.80, width=2.25, height=0.30)

        elif nm == 'TextBox 5':  # "Grid Meter 233" — nudge in from left
            move_and_size(shape, left=0.06, top=1.96, width=2.25, height=0.30)

        elif nm == 'TextBox 96': # "Off-Grid" — nudge in from left
            move_and_size(shape, left=0.06, top=3.58, width=1.80, height=0.32)

        elif nm == 'TextBox 481': # "Yes" — completely off left; place near connector
            move_and_size(shape, left=1.65, top=1.15, width=0.75, height=0.30)

        elif nm == 'TextBox 23': # "8A Backup" — completely off right; place top-center
            move_and_size(shape, left=2.30, top=0.45, width=2.10, height=0.30)

        elif nm == 'TextBox 458': # "Fangzhi Energy Storage" — trim right
            move_and_size(shape, left=6.10, width=1.35, height=2.40)

        elif nm == 'TextBox 524': # "Grid-Connected | Circuit Direction" — trim right
            move_and_size(shape, left=4.55, width=2.88, height=1.40)

        elif nm == 'TextBox 490': # "Off-Grid Circuit Direction" — fine, just set font
            pass

        else:
            clamp(shape)


# Slide 3 (slide_idx=2): startup/shutdown procedures
def fix_slide3(slide):
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        tf = shape.text_frame
        set_font(tf, STANDARD_PT)
        nm = shape.name

        if nm == 'TextBox 19':
            # "Energy Storage ON/OFF" — was above slide; move to left whitespace below pic
            move_and_size(shape, left=0.05, top=4.00, width=3.28, height=0.60)

        elif nm == 'TextBox 29':
            # "(Turn on energy storage first) | Liquid Cooling steps" — above slide
            move_and_size(shape, left=0.05, top=4.68, width=3.28, height=0.90)

        elif nm == 'TextBox 22':  # "Load" — fine, minor tidy
            move_and_size(shape, width=1.25, height=0.30)

        elif nm == 'TextBox 25':  # "Operating Steps" multi-line
            move_and_size(shape, left=3.30, top=7.70, width=3.35, height=1.60)

        elif nm == 'TextBox 27':  # "Shutdown Steps"
            move_and_size(shape, left=5.90, top=7.70, width=1.54, height=0.90)

        elif nm == 'TextBox 28':  # "Startup Steps" (two blocks)
            move_and_size(shape, left=1.52, top=7.82, width=1.85, height=1.05)

        elif nm == 'TextBox 32':  # "(During startup) Check screen…"
            move_and_size(shape, left=3.96, top=6.75, width=3.40, height=0.75)

        else:
            clamp(shape)


# Slide 4 (slide_idx=3): grid-connected vs off-grid mode switching
def fix_slide4(slide):
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        tf = shape.text_frame
        set_font(tf, STANDARD_PT)
        nm = shape.name

        if nm == 'TextBox 17':   # "Step 2: grid switch upper" — trim right
            move_and_size(shape, left=3.94, width=3.49, height=0.50)

        elif nm == 'TextBox 2':  # "Confirm lights / If 3 lights" — trim right
            move_and_size(shape, left=2.55, width=4.88, height=1.15)

        else:
            clamp(shape)


# Slide 5 (slide_idx=4): check mode steps
def fix_slide5(slide):
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        tf = shape.text_frame
        set_font(tf, STANDARD_PT)
        nm = shape.name

        if nm == 'TextBox 11':   # "Password: 88888888" — trim right
            move_and_size(shape, left=4.05, width=3.38, height=0.38)

        else:
            clamp(shape)


# Slide 6 (slide_idx=5): firmware update steps
def fix_slide6(slide):
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        tf = shape.text_frame
        set_font(tf, STANDARD_PT)
        nm = shape.name

        if nm == 'TextBox 23':   # "2) Load .hex…" — trim right
            move_and_size(shape, left=3.90, width=3.53, height=0.38)

        elif nm == 'TextBox 35': # "4) At back of HV box…" — trim right
            move_and_size(shape, left=4.17, width=3.26, height=0.60)

        elif nm == 'TextBox 40': # "3) Connect USB…" — left edge exactly 0, fine
            move_and_size(shape, left=0.05, width=3.68, height=0.38)

        elif nm == 'TextBox 44': # "5) After connecting…"
            move_and_size(shape, left=0.20, width=3.73, height=0.50)

        else:
            clamp(shape)


# ── Run ───────────────────────────────────────────────────────────────────────

FIXERS = {1: fix_slide2, 2: fix_slide3, 3: fix_slide4, 4: fix_slide5, 5: fix_slide6}

for slide_idx in range(1, 6):
    print(f"Fixing slide {slide_idx+1}…")
    FIXERS[slide_idx](prs.slides[slide_idx])

prs.save(dst)
print(f"\nSaved: {dst}")
