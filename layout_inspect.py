import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation

pptx_path = r'C:\Users\JasonOng\OneDrive - Advancer Global Ltd\AST BD\2024 HyESys\Software\Manual logic - bef AI\HZ.SG v2.1 logic2 (EN).pptx'
prs = Presentation(pptx_path)

SW = prs.slide_width
SH = prs.slide_height
print(f'Slide size: {SW.inches:.3f}" x {SH.inches:.3f}" ({int(SW)} x {int(SH)} EMU)')
print()

for slide_idx in range(1, 6):
    slide = prs.slides[slide_idx]
    print(f'=== SLIDE {slide_idx+1} ===')
    for shape in slide.shapes:
        l = shape.left or 0
        t = shape.top or 0
        w = shape.width or 0
        h = shape.height or 0
        l_in = l / 914400
        t_in = t / 914400
        w_in = w / 914400
        h_in = h / 914400
        r_in = (l + w) / 914400
        b_in = (t + h) / 914400
        flags = []
        if (l + w) > SW: flags.append('OUT_RIGHT')
        if (t + h) > SH: flags.append('OUT_BOTTOM')
        if l < 0: flags.append('OUT_LEFT')
        if t < 0: flags.append('OUT_TOP')
        flag_str = ' [' + ','.join(flags) + ']' if flags else ''

        has_text = shape.has_text_frame
        text_preview = ''
        font_sizes = set()
        if has_text:
            text_preview = shape.text_frame.text[:40].replace('\n', ' | ')
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        font_sizes.add(round(run.font.size / 12700))
                if para.runs and not font_sizes:
                    # Check paragraph-level font
                    pass
            # Also check default font at text_frame level

        stype = 'TXT' if has_text else 'IMG'
        fs_str = f' fs={sorted(font_sizes)}' if font_sizes else ''
        print(f'  [{stype}] {shape.name:28s} ({l_in:.2f},{t_in:.2f})-({r_in:.2f},{b_in:.2f}) {w_in:.2f}x{h_in:.2f}{flag_str}{fs_str}  "{text_preview}"')
    print()
