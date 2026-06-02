import sys
import io
import re
import json
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation

def has_chinese(text):
    return bool(re.search(r'[一-鿿＀-￯]', text))

# Load the exact source texts from JSON (extracted directly from the PPTX)
with open(r'C:\Users\JasonOng\AST_Agent\slide_texts.json', encoding='utf-8-sig') as f:
    source_texts = json.load(f)

# Load translations (same order as source_texts, keyed by slide number)
with open(r'C:\Users\JasonOng\AST_Agent\translations.json', encoding='utf-8-sig') as f:
    raw_translations = json.load(f)

# Build lookup dict: exact_chinese_text -> english, per slide
# Use source_texts keys (guaranteed exact) paired with english from translations by index
per_slide_lookup = {}
for slide_num_str, source_list in source_texts.items():
    trans_list = raw_translations.get(slide_num_str, [])
    lookup = {}
    for i, src_text in enumerate(source_list):
        if i < len(trans_list):
            en_text = trans_list[i][1]  # [cn, en] pair -> take en
            lookup[src_text] = en_text
    per_slide_lookup[int(slide_num_str)] = lookup

def apply_translation(para, lookup):
    para_text = para.text
    if not para_text.strip() or not has_chinese(para_text):
        return False

    en_text = lookup.get(para_text)
    if en_text is None:
        # Try substring replacement
        result = para_text
        for cn, en in lookup.items():
            if cn in result:
                result = result.replace(cn, en)
        if result == para_text and has_chinese(result):
            print(f"  [UNTRANSLATED] {repr(para_text)}")
            return False
        en_text = result

    runs = para.runs
    if not runs:
        return False

    print(f"  OK: {repr(para_text[:45])} -> {repr(en_text[:45])}")

    if len(runs) == 1:
        runs[0].text = en_text
    else:
        runs[0].text = en_text
        for run in runs[1:]:
            run.text = ''
    return True

src_path = r'C:\Users\JasonOng\OneDrive - Advancer Global Ltd\AST BD\2024 HyESys\Software\Manual logic - bef AI\HZ.SG v2.1 logic2.pptx'
dst_path = r'C:\Users\JasonOng\OneDrive - Advancer Global Ltd\AST BD\2024 HyESys\Software\Manual logic - bef AI\HZ.SG v2.1 logic2 (EN).pptx'

prs = Presentation(src_path)

print("Translating slides 2-6...\n")
for slide_idx in range(1, 6):
    slide_num = slide_idx + 1
    lookup = per_slide_lookup.get(slide_num, {})
    slide = prs.slides[slide_idx]
    print(f"--- Slide {slide_num} ---")
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                apply_translation(para, lookup)
    print()

prs.save(dst_path)
print(f"Saved: {dst_path}")
